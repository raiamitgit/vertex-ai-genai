import io
import copy
import zipfile
import math
import fitz  # PyMuPDF
from pptx import Presentation

def split_pdf(file_bytes, chunk_size=10):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    total_pages = len(doc)
    chunks = []
    
    for i in range(0, total_pages, chunk_size):
        chunk_doc = fitz.open()
        chunk_doc.insert_pdf(doc, from_page=i, to_page=min(i + chunk_size - 1, total_pages - 1))
        chunks.append(chunk_doc.write())
        chunk_doc.close()
        
    doc.close()
    return chunks

def merge_pdfs(chunks_bytes):
    merged_doc = fitz.open()
    for chunk_bytes in chunks_bytes:
        doc = fitz.open(stream=chunk_bytes, filetype="pdf")
        merged_doc.insert_pdf(doc)
        doc.close()
    result = merged_doc.write()
    merged_doc.close()
    return result

def split_pptx(file_bytes, chunk_size=5):
    prs = Presentation(io.BytesIO(file_bytes))
    total_slides = len(prs.slides)
    chunks = []
    
    for i in range(0, total_slides, chunk_size):
        chunk_prs = Presentation(io.BytesIO(file_bytes))
        xml_slides = chunk_prs.slides._sldIdLst
        slides = list(xml_slides)
        
        for j in range(total_slides - 1, -1, -1):
            if j < i or j >= i + chunk_size:
                xml_slides.remove(slides[j])
                
        out = io.BytesIO()
        chunk_prs.save(out)
        
        # Shrink the PPTX byte array to bypass Vertex AI's 20MB limit
        # By removing embedded videos/images that Vertex doesn't need for XML translation
        chunk_bytes = out.getvalue()
        in_zip = zipfile.ZipFile(io.BytesIO(chunk_bytes), "r")
        shrunk_io = io.BytesIO()
        with zipfile.ZipFile(shrunk_io, "w", zipfile.ZIP_DEFLATED) as out_zip:
            for item in in_zip.infolist():
                if item.filename.startswith("ppt/media/") or item.filename.startswith("ppt/embeddings/"):
                    continue
                data = in_zip.read(item.filename)
                out_zip.writestr(item, data)
        in_zip.close()
        
        chunks.append({
            "bytes": shrunk_io.getvalue(),
            "start_idx": i,
            "end_idx": min(i + chunk_size, total_slides) - 1
        })
        
    return chunks

def merge_pptx(original_bytes, translated_chunks):
    orig_zip = zipfile.ZipFile(io.BytesIO(original_bytes), "r")
    zip_contents = {item.filename: orig_zip.read(item.filename) for item in orig_zip.infolist()}
    orig_zip.close()
    
    for chunk in translated_chunks:
        start = chunk["start_idx"]
        end = chunk["end_idx"]
        chunk_zip = zipfile.ZipFile(io.BytesIO(chunk["bytes"]), "r")
        
        # In the chunk PPTX, the slides are numbered 1 to (end - start + 1)
        for chunk_slide_idx in range(1, (end - start + 1) + 1):
            orig_slide_idx = start + chunk_slide_idx
            
            slide_file = f"ppt/slides/slide{chunk_slide_idx}.xml"
            if slide_file in chunk_zip.namelist():
                orig_slide_file = f"ppt/slides/slide{orig_slide_idx}.xml"
                zip_contents[orig_slide_file] = chunk_zip.read(slide_file)
                
            rels_file = f"ppt/slides/_rels/slide{chunk_slide_idx}.xml.rels"
            if rels_file in chunk_zip.namelist():
                orig_rels_file = f"ppt/slides/_rels/slide{orig_slide_idx}.xml.rels"
                zip_contents[orig_rels_file] = chunk_zip.read(rels_file)
                
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as merged_zip:
        for filename, data in zip_contents.items():
            merged_zip.writestr(filename, data)
            
    return out.getvalue()

def split_docx(file_bytes):
    # DOCX chunking is very complex due to word/document.xml monolithic structure.
    # Passing directly since most DOCX files are < 20MB.
    return [file_bytes]

def merge_docx(chunks_bytes):
    return chunks_bytes[0]
